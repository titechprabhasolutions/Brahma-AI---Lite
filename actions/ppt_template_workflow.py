from __future__ import annotations

import base64
import hashlib
import json
import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, quote_plus, urljoin, urlparse

import requests

try:
    from bs4 import BeautifulSoup
except Exception as exc:  # pragma: no cover - dependency is already bundled, but keep graceful fallback
    BeautifulSoup = None  # type: ignore[assignment]

try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except Exception:  # pragma: no cover
    PP_PLACEHOLDER = None  # type: ignore[assignment]


PROJECT_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_CACHE_DIR = PROJECT_ROOT / "config" / "ppt_template_cache"
TEMPLATE_CACHE_INDEX = TEMPLATE_CACHE_DIR / "index.json"
TEMPLATE_ASSET_DIR = TEMPLATE_CACHE_DIR / "assets"

HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0 Safari/537.36"
    )
}

SUPPORTED_DOWNLOAD_EXTS = {".pptx", ".zip"}
TRUSTED_TEMPLATE_HINTS = (
    "slidesgo",
    "slidescarnival",
    "poweredtemplate",
    "presentationgo",
    "slidebazaar",
    "24slides",
    "templates.office",
    "free-powerpoint-templates",
    "template.net",
    "envato",
)


def _slugify(text: str, default: str = "template") -> str:
    safe = re.sub(r"[^a-zA-Z0-9._ -]+", "", (text or "").strip()).lower()
    safe = re.sub(r"\s+", "-", safe).strip("-._")
    return safe or default


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def _normalize_blob(*parts: Any) -> str:
    blob = " ".join(_clean_text(str(p)) for p in parts if p is not None)
    blob = blob.lower()
    blob = blob.replace("/", " ").replace("\\", " ")
    blob = re.sub(r"[\W_]+", " ", blob)
    return re.sub(r"\s+", " ", blob).strip()


def _unique(seq):
    seen = set()
    out = []
    for item in seq:
        if not item:
            continue
        if item in seen:
            continue
        seen.add(item)
        out.append(item)
    return out


def infer_presentation_profile(
    title: str,
    subtitle: str = "",
    theme_hint: str = "",
    outline: str = "",
    slides: list[dict] | None = None,
) -> dict[str, Any]:
    """Infer a presentation profile and a good search query from the user topic."""

    slide_blob = ""
    if slides:
        slide_blob = " ".join(
            _normalize_blob(
                s.get("title", ""),
                s.get("kicker", ""),
                " ".join(s.get("bullets") or []) if isinstance(s.get("bullets"), list) else s.get("bullets", ""),
                s.get("notes", ""),
            )
            for s in slides
            if isinstance(s, dict)
        )
    blob = _normalize_blob(title, subtitle, theme_hint, outline, slide_blob)

    rules = [
        {
            "name": "cybersecurity",
            "keywords": (
                "cyber security",
                "cybersecurity",
                "security",
                "threat",
                "malware",
                "ransomware",
                "hacker",
                "defense",
                "incident response",
                "penetration",
                "zero trust",
            ),
            "query": "dark cybersecurity technology presentation template",
            "audience": "security professionals",
            "presentation_type": "security briefing",
            "style": "dark technical",
            "palette": "dark red and charcoal",
            "image_terms": ["cyber security", "network security", "digital shield"],
        },
        {
            "name": "technology",
            "keywords": (
                "ai",
                "artificial intelligence",
                "technology",
                "tech",
                "software",
                "cloud",
                "data",
                "machine learning",
                "automation",
                "digital",
                "robotics",
            ),
            "query": "modern futuristic technology powerpoint template",
            "audience": "technology audience",
            "presentation_type": "technology deck",
            "style": "modern futuristic",
            "palette": "neon blue and black",
            "image_terms": ["artificial intelligence", "technology", "futuristic interface"],
        },
        {
            "name": "startup",
            "keywords": (
                "startup",
                "pitch",
                "investor",
                "venture",
                "saas",
                "founder",
                "funding",
                "roadmap",
                "product launch",
                "deck",
            ),
            "query": "startup investor pitch deck powerpoint template",
            "audience": "investors and founders",
            "presentation_type": "pitch deck",
            "style": "clean business",
            "palette": "minimal black and white with red accent",
            "image_terms": ["startup office", "business presentation", "pitch deck"],
        },
        {
            "name": "medical",
            "keywords": (
                "medical",
                "health",
                "healthcare",
                "clinical",
                "hospital",
                "pharma",
                "biotech",
                "medicine",
                "patient",
                "research study",
            ),
            "query": "clean healthcare medical presentation template",
            "audience": "healthcare professionals",
            "presentation_type": "medical presentation",
            "style": "clean clinical",
            "palette": "white blue and teal",
            "image_terms": ["medical healthcare", "clinical research", "hospital technology"],
        },
        {
            "name": "history",
            "keywords": (
                "history",
                "historical",
                "historical",
                "ancient",
                "civilization",
                "heritage",
                "india",
                "indian",
                "mughal",
                "medieval",
                "timeline",
                "war",
                "battle",
                "army",
                "military",
                "patriotic",
                "victory",
                "kargil",
                "vijay diwas",
                "memorial",
            ),
            "query": "classic history education powerpoint template",
            "audience": "students and educators",
            "presentation_type": "educational history deck",
            "style": "classic academic",
            "palette": "warm cream and deep blue",
            "image_terms": ["history timeline", "ancient architecture", "education classroom"],
        },
        {
            "name": "education",
            "keywords": (
                "education",
                "classroom",
                "lesson",
                "student",
                "teacher",
                "study",
                "learning",
                "school",
                "university",
                "research",
            ),
            "query": "clean education presentation template",
            "audience": "students and educators",
            "presentation_type": "education deck",
            "style": "academic minimal",
            "palette": "white and blue",
            "image_terms": ["education", "classroom", "research notes"],
        },
        {
            "name": "corporate",
            "keywords": (
                "corporate",
                "business",
                "board",
                "enterprise",
                "report",
                "strategy",
                "finance",
                "annual",
                "executive",
                "operations",
            ),
            "query": "professional corporate business powerpoint template",
            "audience": "executives and teams",
            "presentation_type": "corporate report",
            "style": "professional minimal",
            "palette": "black white red accent",
            "image_terms": ["corporate strategy", "business analytics", "office meeting"],
        },
        {
            "name": "creative",
            "keywords": (
                "creative",
                "portfolio",
                "marketing",
                "campaign",
                "branding",
                "design",
                "fashion",
                "art",
                "showcase",
            ),
            "query": "creative modern presentation template",
            "audience": "creative teams",
            "presentation_type": "creative presentation",
            "style": "colorful modern",
            "palette": "bold mixed accent colors",
            "image_terms": ["creative design", "marketing presentation", "brand identity"],
        },
        {
            "name": "finance",
            "keywords": (
                "finance",
                "financial",
                "investment",
                "stock",
                "market",
                "revenue",
                "profit",
                "budget",
                "dashboard",
            ),
            "query": "professional finance business presentation template",
            "audience": "finance professionals",
            "presentation_type": "finance deck",
            "style": "corporate minimal",
            "palette": "navy, white, and green",
            "image_terms": ["finance dashboard", "financial analytics", "business growth"],
        },
        {
            "name": "marketing",
            "keywords": (
                "marketing",
                "campaign",
                "social media",
                "brand",
                "advertising",
                "promotion",
                "launch",
                "sales",
            ),
            "query": "modern marketing presentation template",
            "audience": "marketing teams",
            "presentation_type": "marketing deck",
            "style": "modern engaging",
            "palette": "bright modern",
            "image_terms": ["marketing campaign", "social media strategy", "creative launch"],
        },
        {
            "name": "science",
            "keywords": (
                "science",
                "research",
                "experiment",
                "lab",
                "study",
                "analysis",
                "academic paper",
                "biology",
                "chemistry",
                "physics",
            ),
            "query": "clean research presentation template",
            "audience": "researchers and students",
            "presentation_type": "research presentation",
            "style": "academic clean",
            "palette": "blue and white",
            "image_terms": ["scientific research", "laboratory", "data analysis"],
        },
    ]

    selected = None
    selected_score = -1
    for rule in rules:
        score = 0
        for keyword in rule["keywords"]:
            if keyword in blob:
                score += 4
        if score > selected_score:
            selected_score = score
            selected = rule

    if not selected or selected_score <= 0:
        selected = {
            "name": "auto",
            "query": "professional modern powerpoint template",
            "audience": "general audience",
            "presentation_type": "presentation deck",
            "style": "balanced modern",
            "palette": "auto",
            "image_terms": ["presentation template", "modern workspace", "abstract technology"],
        }

    topic = _clean_text(title) or _clean_text(subtitle) or "Presentation"
    if selected["name"] == "technology":
        topic = topic or "Technology"

    primary_query = selected["query"]
    search_queries = _unique(
        [
            f'"{primary_query}"',
            f"{primary_query}",
            f"{primary_query} filetype:pptx",
            f"site:slidesgo.com {primary_query}",
            f"site:slidescarnival.com {primary_query}",
            f"site:poweredtemplate.com {primary_query}",
            f"site:presentationgo.com {primary_query}",
            f"site:slidebazaar.com {primary_query}",
            f"site:templates.office.com {primary_query}",
        ]
    )
    image_terms = _unique([*selected.get("image_terms", []), _clean_text(topic), _clean_text(title), _clean_text(subtitle)])

    return {
        "topic": topic,
        "category": selected["name"],
        "audience": selected["audience"],
        "presentation_type": selected["presentation_type"],
        "style": selected["style"],
        "palette": selected["palette"],
        "search_query": primary_query,
        "search_queries": search_queries,
        "image_queries": image_terms,
        "cache_key": hashlib.sha1(primary_query.encode("utf-8")).hexdigest()[:16],
    }


def _load_cache_index() -> dict[str, Any]:
    if TEMPLATE_CACHE_INDEX.exists():
        try:
            return json.loads(TEMPLATE_CACHE_INDEX.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _save_cache_index(index: dict[str, Any]) -> None:
    TEMPLATE_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    TEMPLATE_CACHE_INDEX.write_text(json.dumps(index, indent=2), encoding="utf-8")


def _unwrap_bing_url(href: str) -> str:
    if not href:
        return href
    parsed = urlparse(href)
    if "bing.com" not in parsed.netloc.lower():
        return href
    qs = parse_qs(parsed.query)
    payload = qs.get("u", [None])[0]
    if payload and payload.startswith("a1"):
        payload = payload[2:]
        payload += "=" * (-len(payload) % 4)
        try:
            return base64.b64decode(payload).decode("utf-8", "ignore")
        except Exception:
            return href
    return href


def _is_downloadable_template_url(url: str) -> bool:
    lower = url.lower().split("?", 1)[0]
    return any(lower.endswith(ext) for ext in SUPPORTED_DOWNLOAD_EXTS)


def _search_bing(query: str, max_results: int = 8) -> list[dict[str, str]]:
    try:
        resp = requests.get(
            "https://www.bing.com/search",
            params={"q": query, "count": max_results, "setlang": "en-US", "cc": "us"},
            headers=HEADERS,
            timeout=20,
        )
        resp.raise_for_status()
    except Exception:
        return []

    if BeautifulSoup is None:
        return []

    soup = BeautifulSoup(resp.text, "html.parser")
    results = []
    for li in soup.select("li.b_algo"):
        a = li.select_one("h2 a")
        if not a:
            continue
        href = _unwrap_bing_url(a.get("href", "") or "")
        title = _clean_text(a.get_text(" ", strip=True))
        snippet = _clean_text(li.get_text(" ", strip=True))
        if not href or not title:
            continue
        results.append({"title": title, "href": href, "snippet": snippet})
    return results


def _presentationgo_category_urls(profile: dict[str, Any]) -> list[str]:
    category = profile.get("category", "auto")
    base = "https://www.presentationgo.com/presentation/category/templates"
    mapping = {
        "technology": [f"{base}/technology/", f"{base}/abstract/"],
        "startup": [f"{base}/business/", f"{base}/abstract/"],
        "corporate": [f"{base}/business/", f"{base}/abstract/"],
        "finance": [f"{base}/business/", f"{base}/abstract/"],
        "medical": [f"{base}/health-recreation/", f"{base}/science/"],
        "science": [f"{base}/science/", f"{base}/abstract/"],
        "education": [f"{base}/education/", f"{base}/abstract/"],
        "history": [f"{base}/education/", f"{base}/religion-holidays/", f"{base}/abstract/"],
        "cybersecurity": [f"{base}/technology/", f"{base}/abstract/"],
        "creative": [f"{base}/abstract/", f"{base}/color-palettes/"],
        "marketing": [f"{base}/business/", f"{base}/abstract/"],
        "auto": [f"{base}/abstract/", f"{base}/business/", f"{base}/technology/"],
    }
    return mapping.get(category, mapping["auto"])


def _fetch_html(url: str) -> str | None:
    try:
        resp = requests.get(url, headers=HEADERS, timeout=25)
        resp.raise_for_status()
        return resp.text
    except Exception:
        return None


def _presentationgo_category_articles(category_url: str) -> list[dict[str, str]]:
    html = _fetch_html(category_url)
    if not html:
        return []
    if BeautifulSoup is None:
        return []
    soup = BeautifulSoup(html, "html.parser")
    seen = set()
    articles = []
    for a in soup.select('a[href*="/presentation/"]'):
        href = a.get("href") or ""
        if "/presentation/category/" in href or "/presentation/tag/" in href or "/download/" in href:
            continue
        href = urljoin(category_url, href)
        href = href.split("#", 1)[0]
        if href in seen:
            continue
        seen.add(href)
        text = _clean_text(a.get_text(" ", strip=True))
        if not text:
            text = Path(urlparse(href).path.rstrip("/")).name.replace("-", " ")
        articles.append({"title": text, "href": href})
    return articles


def _presentationgo_download_links(article_url: str) -> list[str]:
    html = _fetch_html(article_url)
    if not html:
        return []
    if BeautifulSoup is None:
        return []
    soup = BeautifulSoup(html, "html.parser")
    links = []
    for a in soup.select('a[href]'):
        href = a.get("href") or ""
        text = _clean_text(a.get_text(" ", strip=True)).lower()
        absolute = urljoin(article_url, href)
        absolute_lower = absolute.lower()
        if "presentationgo.com/download/" in absolute_lower and ("potx" in text or "widescreen" in text or "ppt" in text):
            links.append(absolute)
        elif _is_downloadable_template_url(absolute_lower):
            links.append(absolute)
    return _unique(links)


def _presentationgo_choose_article(profile: dict[str, Any]) -> tuple[str, str] | None:
    best = None
    best_score = -1
    query = _normalize_blob(profile.get("search_query", ""), profile.get("presentation_type", ""), profile.get("style", ""))
    for category_url in _presentationgo_category_urls(profile):
        articles = _presentationgo_category_articles(category_url)
        for article in articles:
            title = (article.get("title") or "").lower()
            href = (article.get("href") or "").lower()
            score = 0
            if profile["category"] in title or profile["presentation_type"] in title:
                score += 12
            if profile["style"] in title:
                score += 8
            if any(term in title for term in query.split()):
                score += 4
            if any(term in href for term in ("template", "presentation")):
                score += 2
            if "history" in query and any(term in title for term in ("history", "war", "india", "army", "patriotic", "vijay", "independence", "freedom")):
                score += 8
            if "technology" in query and any(term in title for term in ("tech", "digital", "future", "data", "aerospace", "blue")):
                score += 8
            if "startup" in query and any(term in title for term in ("business", "strategy", "executive", "momentum", "pitch")):
                score += 8
            if score > best_score:
                best_score = score
                best = article
    if best:
        return best.get("title", ""), best.get("href", "")
    return None


def _candidate_score(candidate: dict[str, str], profile: dict[str, Any], query: str) -> int:
    title = (candidate.get("title") or "").lower()
    href = (candidate.get("href") or "").lower()
    snippet = (candidate.get("snippet") or "").lower()
    score = 0

    if _is_downloadable_template_url(href):
        score += 35
    if any(h in href for h in TRUSTED_TEMPLATE_HINTS):
        score += 20
    if "template" in title or "template" in snippet or "presentation" in title or "presentation" in snippet:
        score += 10
    if profile["category"] in title or profile["category"] in snippet:
        score += 10
    if profile["presentation_type"] in title or profile["presentation_type"] in snippet:
        score += 6
    if profile["style"] in title or profile["style"] in snippet:
        score += 6
    if profile["audience"] in title or profile["audience"] in snippet:
        score += 4
    query_terms = [term for term in query.lower().split() if len(term) > 3]
    hits = sum(1 for term in query_terms if term in title or term in snippet or term in href)
    score += min(hits * 2, 12)
    if any(bad in href for bad in ("dictionary", "wiki", "how-to", "tutorial", "example", "definition")):
        score -= 6
    return score


def _extract_page_download_links(page_url: str, profile: dict[str, Any]) -> list[str]:
    try:
        resp = requests.get(page_url, headers=HEADERS, timeout=20)
        resp.raise_for_status()
    except Exception:
        return []
    if BeautifulSoup is None:
        return []

    soup = BeautifulSoup(resp.text, "html.parser")
    links: list[str] = []

    for tag in soup.select("a[href]"):
        href = tag.get("href") or ""
        text = _clean_text(tag.get_text(" ", strip=True)).lower()
        absolute = urljoin(page_url, href)
        absolute_lower = absolute.lower()
        if _is_downloadable_template_url(absolute_lower):
            links.append(absolute)
            continue
        if any(token in absolute_lower for token in ("download", "pptx", "presentation", "template")) and any(
            token in text for token in ("download", "ppt", "pptx", "template")
        ):
            links.append(absolute)

    # look for JSON/inline data blocks that expose a direct download
    for pattern in (r"https?://[^\s\"'<>]+\.pptx", r"https?://[^\s\"'<>]+\.zip"):
        for match in re.findall(pattern, resp.text, flags=re.I):
            links.append(match)

    unique_links = []
    seen = set()
    for link in links:
        if link in seen:
            continue
        seen.add(link)
        unique_links.append(link)
    return unique_links


def _validate_pptx(path: Path) -> bool:
    if not path.exists() or path.stat().st_size == 0:
        return False
    try:
        if zipfile.is_zipfile(path):
            with zipfile.ZipFile(path, "r") as zf:
                names = set(zf.namelist())
                if "[Content_Types].xml" in names and any(name.startswith("ppt/") for name in names):
                    return True
        return path.suffix.lower() == ".pptx" and path.stat().st_size > 0
    except Exception:
        return False


def _materialize_template_file(downloaded_path: Path, dest_path: Path) -> Path | None:
    if _validate_pptx(downloaded_path):
        shutil.copy2(downloaded_path, dest_path)
        return dest_path

    if zipfile.is_zipfile(downloaded_path):
        with tempfile.TemporaryDirectory(prefix="brahma_pptx_extract_") as tmpdir:
            extract_dir = Path(tmpdir)
            with zipfile.ZipFile(downloaded_path, "r") as zf:
                zf.extractall(extract_dir)
            for candidate in extract_dir.rglob("*.pptx"):
                if _validate_pptx(candidate):
                    shutil.copy2(candidate, dest_path)
                    return dest_path
    return None


def _download_url(url: str, dest: Path) -> Path | None:
    try:
        with requests.get(url, headers=HEADERS, timeout=40, stream=True, allow_redirects=True) as resp:
            resp.raise_for_status()
            tmp = dest.with_suffix(".download")
            with open(tmp, "wb") as fh:
                for chunk in resp.iter_content(chunk_size=1024 * 64):
                    if chunk:
                        fh.write(chunk)
            return _materialize_template_file(tmp, dest)
    except Exception:
        return None
    finally:
        try:
            if dest.with_suffix(".download").exists():
                dest.with_suffix(".download").unlink()
        except Exception:
            pass


def _choose_template_candidate(results: list[dict[str, str]], profile: dict[str, Any], query: str) -> dict[str, str] | None:
    scored = sorted(results, key=lambda item: _candidate_score(item, profile, query), reverse=True)
    for candidate in scored:
        href = candidate.get("href") or ""
        if _is_downloadable_template_url(href):
            return candidate
        download_links = _extract_page_download_links(href, profile)
        if download_links:
            return {**candidate, "download_url": download_links[0]}
    return None


def resolve_presentation_template(profile: dict[str, Any]) -> dict[str, Any] | None:
    TEMPLATE_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    TEMPLATE_ASSET_DIR.mkdir(parents=True, exist_ok=True)
    cache_index = _load_cache_index()
    cache_key = profile["cache_key"]

    cached = cache_index.get(cache_key)
    if cached:
        cached_path = Path(cached.get("path", ""))
        if cached_path.exists() and _validate_pptx(cached_path):
            return {
                "path": cached_path,
                "source_url": cached.get("source_url"),
                "search_query": cached.get("search_query"),
                "cached": True,
                "profile": profile,
            }

    chosen = _presentationgo_choose_article(profile)
    if chosen:
        _, article_url = chosen
        download_links = _presentationgo_download_links(article_url)
        for download_url in download_links:
            dest_name = f"{cache_key}.pptx"
            dest_path = TEMPLATE_ASSET_DIR / dest_name
            downloaded = _download_url(download_url, dest_path)
            if downloaded and _validate_pptx(downloaded):
                cache_index[cache_key] = {
                    "path": str(downloaded),
                    "source_url": download_url,
                    "search_query": profile.get("search_query"),
                    "cached_at": __import__("time").time(),
                    "profile": profile,
                }
                _save_cache_index(cache_index)
                return {
                    "path": downloaded,
                    "source_url": download_url,
                    "search_query": profile.get("search_query"),
                    "cached": False,
                    "profile": profile,
                }

    for query in profile.get("search_queries", []):
        results = _search_bing(query, max_results=8)
        if not results:
            continue
        candidate = _choose_template_candidate(results, profile, query)
        if not candidate:
            continue
        download_url = candidate.get("download_url") or candidate.get("href")
        if not download_url:
            continue
        dest_name = f"{cache_key}.pptx"
        dest_path = TEMPLATE_ASSET_DIR / dest_name
        downloaded = _download_url(download_url, dest_path)
        if downloaded and _validate_pptx(downloaded):
            cache_index[cache_key] = {
                "path": str(downloaded),
                "source_url": download_url,
                "search_query": query,
                "cached_at": __import__("time").time(),
                "profile": profile,
            }
            _save_cache_index(cache_index)
            return {
                "path": downloaded,
                "source_url": download_url,
                "search_query": query,
                "cached": False,
                "profile": profile,
            }

    return None


def _clear_presentation(prs) -> None:
    slide_id_list = prs.slides._sldIdLst
    for idx in range(len(prs.slides) - 1, -1, -1):
        rel = slide_id_list[idx].rId
        prs.part.drop_rel(rel)
        del slide_id_list[idx]


def _placeholder_type_name(ph_type: Any) -> str:
    if PP_PLACEHOLDER is None:
        return str(ph_type).lower()
    mapping = {
        PP_PLACEHOLDER.TITLE: "title",
        PP_PLACEHOLDER.CENTER_TITLE: "title",
        PP_PLACEHOLDER.SUBTITLE: "subtitle",
        PP_PLACEHOLDER.BODY: "body",
        PP_PLACEHOLDER.OBJECT: "body",
        PP_PLACEHOLDER.PICTURE: "picture",
        PP_PLACEHOLDER.CHART: "chart",
        PP_PLACEHOLDER.TABLE: "table",
        PP_PLACEHOLDER.SLIDE_NUMBER: "meta",
        PP_PLACEHOLDER.DATE: "meta",
        PP_PLACEHOLDER.FOOTER: "meta",
    }
    return mapping.get(ph_type, str(ph_type).lower())


def _iter_placeholders(slide):
    placeholders = []
    for shape in slide.placeholders:
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        placeholders.append((shape, _placeholder_type_name(ph_type)))
    return placeholders


def _fill_text_shape(shape, text: str, bullets: list[str] | None = None, font_size: int | None = None) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    items = bullets or [text]
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = _clean_text(str(item))
        p.level = 0
        if font_size:
            for run in p.runs:
                run.font.size = font_size


def _pick_layout(prs, has_image: bool, is_title: bool):
    best_layout = None
    best_score = -1
    for layout in prs.slide_layouts:
        layout_types = []
        for ph in layout.placeholders:
            try:
                layout_types.append(_placeholder_type_name(ph.placeholder_format.type))
            except Exception:
                pass
        score = 0
        if is_title:
            if "title" in layout_types:
                score += 10
            if "subtitle" in layout_types:
                score += 6
            if "body" in layout_types:
                score += 2
        else:
            if "title" in layout_types:
                score += 8
            if "body" in layout_types:
                score += 8
            if has_image and "picture" in layout_types:
                score += 10
            if not layout_types:
                score -= 3
        if score > best_score:
            best_score = score
            best_layout = layout
    return best_layout or prs.slide_layouts[0]


def _load_related_images(profile: dict[str, Any], slides: list[dict[str, Any]], max_images: int = 6) -> list[Path]:
    try:
        requests.get("https://commons.wikimedia.org", timeout=5, headers=HEADERS)
    except Exception:
        return []

    image_cache_dir = TEMPLATE_ASSET_DIR / "images" / profile["cache_key"]
    image_cache_dir.mkdir(parents=True, exist_ok=True)
    queries = _unique(
        [
            *profile.get("image_queries", []),
            *[
                _clean_text(s.get("title", ""))
                for s in slides
                if isinstance(s, dict) and s.get("title")
            ],
        ]
    )

    results: list[Path] = []
    for query in queries:
        if len(results) >= max_images:
            break
        try:
            resp = requests.get(
                "https://commons.wikimedia.org/w/api.php",
                params={
                    "action": "query",
                    "generator": "search",
                    "gsrsearch": query,
                    "gsrnamespace": 6,
                    "gsrlimit": 5,
                    "prop": "imageinfo",
                    "iiprop": "url",
                    "format": "json",
                },
                headers=HEADERS,
                timeout=20,
            )
            resp.raise_for_status()
            payload = resp.json()
        except Exception:
            continue

        pages = payload.get("query", {}).get("pages", {})
        candidates = []
        for page in pages.values():
            infos = page.get("imageinfo") or []
            if not infos:
                continue
            image_url = infos[0].get("url")
            if image_url:
                candidates.append(image_url)

        for image_url in candidates:
            if len(results) >= max_images:
                break
            ext = Path(urlparse(image_url).path).suffix.lower() or ".jpg"
            dest = image_cache_dir / f"{hashlib.sha1(image_url.encode('utf-8')).hexdigest()[:16]}{ext}"
            if dest.exists() and dest.stat().st_size > 0:
                results.append(dest)
                continue
            try:
                with requests.get(image_url, headers=HEADERS, timeout=25, stream=True) as resp_img:
                    resp_img.raise_for_status()
                    with open(dest, "wb") as fh:
                        for chunk in resp_img.iter_content(chunk_size=1024 * 64):
                            if chunk:
                                fh.write(chunk)
                if dest.exists() and dest.stat().st_size > 0:
                    results.append(dest)
            except Exception:
                try:
                    if dest.exists():
                        dest.unlink()
                except Exception:
                    pass
                continue
    return results


def build_presentation_from_template(
    template_path: Path,
    parameters: dict,
    profile: dict[str, Any],
    slides: list[dict[str, Any]],
    output_path: Path,
    auto_open: bool = True,
    player=None,
) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore

    prs = Presentation(str(template_path))
    _clear_presentation(prs)
    image_pool = _load_related_images(profile, slides)
    image_idx = 0

    def _next_image() -> Path | None:
        nonlocal image_idx
        if image_idx >= len(image_pool):
            return None
        img = image_pool[image_idx]
        image_idx += 1
        return img

    def _title_text(spec: dict[str, Any], idx: int) -> str:
        title = _clean_text(spec.get("title", f"Slide {idx}"))
        kicker = _clean_text(spec.get("kicker", f"Slide {idx:02d}"))
        return title or kicker

    for idx, spec in enumerate(slides, 1):
        has_image = bool(image_pool)
        is_title = idx == 1
        layout = _pick_layout(prs, has_image=has_image, is_title=is_title)
        slide = prs.slides.add_slide(layout)
        placeholders = _iter_placeholders(slide)
        title_text = _title_text(spec, idx)
        bullet_items = []
        raw_bullets = spec.get("bullets") or []
        for bullet in raw_bullets:
            if isinstance(bullet, dict):
                bullet_items.append(_clean_text(bullet.get("text", "")))
            else:
                bullet_items.append(_clean_text(str(bullet)))
        if not bullet_items:
            notes = _clean_text(spec.get("notes", ""))
            if notes:
                bullet_items = [notes]
        if not bullet_items:
            bullet_items = [profile.get("presentation_type", "Key idea"), profile.get("audience", "Audience")]

        picture_shape = None
        title_shape = None
        subtitle_shape = None
        body_shapes = []
        for shape, ph_type in placeholders:
            if ph_type == "picture" and picture_shape is None:
                picture_shape = shape
            elif ph_type == "title" and title_shape is None:
                title_shape = shape
            elif ph_type == "subtitle" and subtitle_shape is None:
                subtitle_shape = shape
            elif ph_type in ("body", "content", "object", "table", "chart"):
                body_shapes.append(shape)

        if title_shape is not None:
            try:
                _fill_text_shape(title_shape, title_text)
            except Exception:
                pass

        if subtitle_shape is not None:
            sub_text = _clean_text(spec.get("notes") or profile.get("audience") or profile.get("style") or "")
            if not sub_text:
                sub_text = profile.get("presentation_type", "")
            try:
                _fill_text_shape(subtitle_shape, sub_text)
            except Exception:
                pass

        body_text = " • ".join(bullet_items[:3]) if bullet_items else title_text
        if body_shapes:
            first_body = body_shapes[0]
            try:
                _fill_text_shape(first_body, body_text, bullets=bullet_items)
            except Exception:
                try:
                    first_body.text = body_text
                except Exception:
                    pass

        if picture_shape is not None:
            image_path = _next_image()
            if image_path and image_path.exists():
                try:
                    picture_shape.insert_picture(str(image_path))
                except Exception:
                    try:
                        left = picture_shape.left
                        top = picture_shape.top
                        width = picture_shape.width
                        height = picture_shape.height
                        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
                    except Exception:
                        pass

        # If the template layout has no useful body placeholder, add a minimal text box so content is not lost.
        if not body_shapes and not picture_shape:
            try:
                tx = slide.shapes.add_textbox(914400, 1828800, 7620000, 2286000)
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = body_text
            except Exception:
                pass

    prs.save(str(output_path))
    if auto_open:
        try:
            from actions.office_builder import _open_file  # local import to avoid circular hint in type checkers

            _open_file(output_path)
        except Exception:
            pass
    return f"Presentation created: {output_path}"
