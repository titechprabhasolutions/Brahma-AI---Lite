"""
office_builder.py - Brahma AI office document generation

Creates PowerPoint presentations and Excel workbooks from structured inputs.
"""

from __future__ import annotations

import json
import os
import subprocess
import re
import sys
from pathlib import Path

from actions.ppt_template_workflow import (
    build_presentation_from_template,
    infer_presentation_profile,
    resolve_presentation_template,
)


PROJECT_NAME = "Brahma AI - Lite"
DEFAULT_OUTPUT_DIR = Path.home() / "Desktop" / "BrahmaAI"


def _sanitize_filename(name: str, default: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9._ -]+", "", (name or "").strip())
    safe = re.sub(r"\s+", " ", safe).strip().replace(" ", "_")
    return safe or default


def _resolve_output_path(output_path: str | None, title: str, ext: str) -> Path:
    if output_path:
        path = Path(output_path).expanduser()
        if not path.is_absolute():
            head = path.parts[0].lower() if path.parts else ""
            tail = Path(*path.parts[1:]) if len(path.parts) > 1 else Path(path.name)
            if head in {"downloads", "download"}:
                path = Path.home() / "Downloads" / tail
            elif head == "desktop":
                path = Path.home() / "Desktop" / tail
            else:
                path = Path.cwd() / path
        if path.suffix.lower() != ext:
            path = path.with_suffix(ext)
        path.parent.mkdir(parents=True, exist_ok=True)
        return path

    DEFAULT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return DEFAULT_OUTPUT_DIR / f"{_sanitize_filename(title, 'brahma_ai_output')}{ext}"


def _parse_json_arg(value, fallback):
    if value is None:
        return fallback
    if isinstance(value, (dict, list)):
        return value
    if isinstance(value, str):
        txt = value.strip()
        if not txt:
            return fallback
        try:
            return json.loads(txt)
        except Exception:
            return fallback
    return fallback


def _coerce_cell_value(value):
    if not isinstance(value, str):
        return value
    text = value.strip()
    if not text:
        return ""
    if text.startswith("="):
        return text
    if text.lower() in {"true", "false"}:
        return text.lower() == "true"
    if re.fullmatch(r"-?\d+", text):
        try:
            return int(text)
        except Exception:
            return text
    if re.fullmatch(r"-?\d+\.\d+", text):
        try:
            return float(text)
        except Exception:
            return text
    return value


def _open_file(path: Path) -> None:
    try:
        if os.name == "nt":
            os.startfile(str(path))  # type: ignore[attr-defined]
            return
        if sys.platform == "darwin":
            subprocess.Popen(["open", str(path)])
            return
        subprocess.Popen(["xdg-open", str(path)])
    except Exception:
        pass


def _import_pptx():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        return Presentation, RGBColor, MSO_SHAPE, Inches, Pt
    except Exception as e:
        raise RuntimeError(
            "python-pptx is required. Install it with: pip install python-pptx"
        ) from e


def _import_openpyxl():
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter
        return Workbook, BarChart, LineChart, PieChart, Reference, Alignment, Font, PatternFill, get_column_letter
    except Exception as e:
        raise RuntimeError(
            "openpyxl is required. Install it with: pip install openpyxl"
        ) from e


def _office_theme():
    return {
        "bg": "07131C",
        "panel": "0E2230",
        "accent": "00D4FF",
        "accent2": "FF8A3D",
        "text": "F1FAFF",
        "muted": "8DB7C8",
        "line": "23485E",
    }


def _theme_library():
    return {
        "auto": {
            "bg": "07131C",
            "panel": "0E2230",
            "accent": "00D4FF",
            "accent2": "FF8A3D",
            "accent3": "7C5CFF",
            "text": "F1FAFF",
            "muted": "8DB7C8",
            "line": "23485E",
            "glow": "103549",
            "surface": "0B1D28",
        },
        "neon": {
            "bg": "05070C",
            "panel": "111827",
            "accent": "21E6C1",
            "accent2": "7C5CFF",
            "accent3": "00D4FF",
            "text": "F6FAFF",
            "muted": "98A9C0",
            "line": "253245",
            "glow": "0D2533",
            "surface": "0B1220",
        },
        "corporate": {
            "bg": "081018",
            "panel": "102130",
            "accent": "F97316",
            "accent2": "22C55E",
            "accent3": "38BDF8",
            "text": "F8FAFC",
            "muted": "94A3B8",
            "line": "2B4057",
            "glow": "132433",
            "surface": "0F172A",
        },
        "luxury": {
            "bg": "0A0910",
            "panel": "161320",
            "accent": "D4AF37",
            "accent2": "F3E8C7",
            "accent3": "A855F7",
            "text": "FFFDF7",
            "muted": "C9C1B2",
            "line": "3B314E",
            "glow": "191224",
            "surface": "11111A",
        },
        "academic": {
            "bg": "0D1117",
            "panel": "141A22",
            "accent": "60A5FA",
            "accent2": "F59E0B",
            "accent3": "34D399",
            "text": "F8FAFC",
            "muted": "94A3B8",
            "line": "263445",
            "glow": "102235",
            "surface": "0F1720",
        },
        "sunset": {
            "bg": "1A0F14",
            "panel": "25141A",
            "accent": "FB7185",
            "accent2": "FDBA74",
            "accent3": "A78BFA",
            "text": "FFF7F8",
            "muted": "E5B7C0",
            "line": "402430",
            "glow": "2A131E",
            "surface": "1E1016",
        },
    }


def _normalize_theme(theme_text: str | None, title: str = "", subtitle: str = "") -> str:
    text = f"{theme_text or ''} {title} {subtitle}".lower()
    if any(word in text for word in ["neon", "futur", "tech", "cyber", "ai", "startup", "saas", "studio", "dash"]):
        return "neon"
    if any(word in text for word in ["luxury", "premium", "gold", "fashion", "brand"]):
        return "luxury"
    if any(word in text for word in ["finance", "corp", "board", "enterprise", "business"]):
        return "corporate"
    if any(word in text for word in ["academic", "research", "science", "education", "study"]):
        return "academic"
    if any(word in text for word in ["sunset", "creative", "marketing", "campaign", "portfolio"]):
        return "sunset"
    return "auto"


def _theme_pack(theme_text: str | None, title: str = "", subtitle: str = "") -> dict:
    theme_key = _normalize_theme(theme_text, title, subtitle)
    return _theme_library().get(theme_key, _theme_library()["auto"])


def _fill_shape(shape, color_hex: str, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color_hex)
    shape.line.fill.background()


def _add_bg(slide, theme, rgb, MSO_SHAPE, Inches, variant: str = "content"):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    _fill_shape(bg, theme["bg"], rgb)

    glow_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(-0.8), Inches(4.5), Inches(4.5))
    _fill_shape(glow_left, theme["glow"], rgb)
    glow_left.fill.transparency = 0.35

    glow_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(4.7), Inches(3.8), Inches(3.8))
    _fill_shape(glow_right, theme["accent"], rgb)
    glow_right.fill.transparency = 0.78

    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22))
    _fill_shape(rail, theme["accent"], rgb)

    if variant == "cover":
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.25), Inches(0.55), Inches(3.55), Inches(6.35))
        _fill_shape(panel, theme["panel"], rgb)
        panel.fill.transparency = 0.08
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.25), Inches(0.55), Inches(0.12), Inches(6.35))
        _fill_shape(accent_bar, theme["accent2"], rgb)
    elif variant == "section":
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.1))
        _fill_shape(band, theme["line"], rgb)
        band.fill.transparency = 0.25
    else:
        left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.32), Inches(0.65), Inches(0.08), Inches(5.9))
        _fill_shape(left, theme["accent2"], rgb)
        right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.95), Inches(0.65), Inches(0.06), Inches(5.9))
        _fill_shape(right, theme["accent3"], rgb)


def _content_layout_variant(idx: int, bullets: list) -> str:
    if idx == 0:
        return "cover"
    if idx % 5 == 0 or (bullets and len(bullets) <= 2):
        return "section"
    if idx % 2 == 0:
        return "split"
    return "content"


def _add_slide_title(slide, title: str, kicker: str, rgb, pt, Inches, MSO_SHAPE):
    colors = _office_theme()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(colors["bg"])
    bg.line.fill.background()

    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.45), Inches(0.14), Inches(0.14))
    marker.fill.solid()
    marker.fill.fore_color.rgb = rgb(colors["accent"])
    marker.line.fill.background()

    kicker_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.38), Inches(2.5), Inches(0.28))
    tf = kicker_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (kicker or PROJECT_NAME).upper()
    r.font.name = "Aptos"
    r.font.size = pt(10)
    r.font.bold = True
    r.font.color.rgb = rgb(colors["accent"])

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(10.7), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = pt(28)
    r.font.bold = True
    r.font.color.rgb = rgb(colors["text"])

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.65), Inches(3.0), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(colors["accent2"])
    rule.line.fill.background()


def _slides_from_outline(outline: str | None, provided: list[dict] | None) -> list[dict]:
    if provided:
        return provided

    if not outline:
        return [
            {"title": "Overview", "bullets": ["Add a slide outline or structured slides input."]},
        ]

    blocks = [b.strip() for b in re.split(r"\n\s*\n", outline.strip()) if b.strip()]
    slides = []
    for idx, block in enumerate(blocks, 1):
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        title = lines[0].lstrip("#").strip() if lines else f"Slide {idx}"
        bullets = []
        for ln in lines[1:]:
            bullet = re.sub(r"^[-*•\d.)\s]+", "", ln).strip()
            if bullet:
                bullets.append(bullet)
        if not bullets and len(lines) > 1:
            bullets = lines[1:]
        slides.append({"title": title, "bullets": bullets or ["Add supporting points here."]})
    return slides


def create_presentation(parameters: dict, player=None) -> str:
    Presentation, RGBColor, MSO_SHAPE, Inches, Pt = _import_pptx()
    rgb = lambda hex_str: RGBColor.from_string(hex_str)

    title = (parameters.get("title") or parameters.get("topic") or PROJECT_NAME).strip()
    subtitle = (parameters.get("subtitle") or parameters.get("audience") or "").strip()
    theme_hint = parameters.get("theme") or parameters.get("visual_theme") or parameters.get("style")
    theme = _theme_pack(theme_hint, title, subtitle)
    output_path = _resolve_output_path(parameters.get("output_path"), title, ".pptx")
    auto_open = parameters.get("auto_open", True)
    slides = _slides_from_outline(parameters.get("outline"), _parse_json_arg(parameters.get("slides"), None))
    slides = slides[:20] if slides else slides

    profile = infer_presentation_profile(title, subtitle=subtitle, theme_hint=theme_hint or "", outline=parameters.get("outline") or "", slides=slides)
    template_choice = resolve_presentation_template(profile)
    if template_choice:
        try:
            return build_presentation_from_template(
                template_path=template_choice["path"],
                parameters=parameters,
                profile=profile,
                slides=slides,
                output_path=output_path,
                auto_open=auto_open,
                player=player,
            )
        except Exception as exc:
            if player and hasattr(player, "add_system_message"):
                try:
                    player.add_system_message(
                        f"Template-based presentation failed, falling back to the built-in designer: {exc}"
                    )
                except Exception:
                    pass

    def add_textbox(slide, left, top, width, height, text, font_size=18, color_key="text", bold=False,
                    font_name="Aptos", align=None, italic=False, all_caps=False):
        color_hex = theme.get(color_key, color_key)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = text.upper() if all_caps else text
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = rgb(color_hex)
        return box

    def add_bullet_card(slide, left, top, width, height, heading, bullets, accent_key="accent", dense=False):
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        _fill_shape(panel, theme["panel"], rgb)
        panel.fill.transparency = 0.06
        edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
        _fill_shape(edge, theme[accent_key], rgb)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
                    heading, font_size=14, color_key=accent_key, bold=True, all_caps=True)
        body = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.52), width - Inches(0.35), height - Inches(0.6))
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        for i, bullet in enumerate(bullets[:6] if dense else bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            txt = bullet.get("text", "") if isinstance(bullet, dict) else str(bullet)
            p.text = f"- {txt}"
            p.level = 0
            p.space_after = Pt(4)
            for run in p.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(15 if dense else 17)
                run.font.color.rgb = rgb(theme["text"])
        return panel

    def add_kpi(slide, left, top, width, height, label, value, accent_key="accent"):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        _fill_shape(card, theme["surface"], rgb)
        card.fill.transparency = 0.04
        outline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
        _fill_shape(outline, theme[accent_key], rgb)
        outline.fill.transparency = 0.18
        add_textbox(slide, left + Inches(0.16), top + Inches(0.12), width - Inches(0.3), Inches(0.22),
                    label, font_size=10, color_key="muted", bold=True, all_caps=True)
        add_textbox(slide, left + Inches(0.16), top + Inches(0.34), width - Inches(0.3), Inches(0.42),
                    value, font_size=22, color_key="text", bold=True)

    def add_section_header(slide, title_text, kicker_text, left, top, width, accent_key="accent"):
        add_textbox(slide, left, top, width, Inches(0.3), kicker_text, font_size=10, color_key=accent_key, bold=True, all_caps=True)
        add_textbox(slide, left, top + Inches(0.22), width, Inches(0.55), title_text, font_size=24, color_key="text", bold=True)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.83), Inches(2.0), Inches(0.06))
        _fill_shape(rule, theme[accent_key], rgb)
        rule.fill.transparency = 0.1

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme, rgb, MSO_SHAPE, Inches, variant="cover")

    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(3.2), Inches(0.25), PROJECT_NAME, font_size=10, color_key="muted", bold=True, all_caps=True)
    add_textbox(slide, Inches(0.7), Inches(0.82), Inches(7.4), Inches(0.75), title, font_size=28, color_key="text", bold=True, font_name="Aptos Display")
    if subtitle:
        add_textbox(slide, Inches(0.72), Inches(1.62), Inches(6.8), Inches(0.45), subtitle, font_size=14, color_key="muted")

    title_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(0.85), Inches(2.85), Inches(5.95))
    _fill_shape(title_panel, theme["panel"], rgb)
    title_panel.fill.transparency = 0.08
    add_textbox(slide, Inches(9.8), Inches(1.08), Inches(2.2), Inches(0.25), "Deck Theme", font_size=10, color_key="muted", bold=True, all_caps=True)
    add_textbox(slide, Inches(9.78), Inches(1.35), Inches(2.25), Inches(1.05), theme_hint or "Auto", font_size=22, color_key="accent", bold=True)
    add_textbox(slide, Inches(9.78), Inches(2.52), Inches(2.1), Inches(1.4),
                "Polished backgrounds\nSection breaks\nSplit layouts\nVisual callouts",
                font_size=13, color_key="text")
    add_kpi(slide, Inches(0.7), Inches(5.05), Inches(1.9), Inches(0.8), "Slides", f"{max(len(slides), 1):02d}", accent_key="accent")
    add_kpi(slide, Inches(2.75), Inches(5.05), Inches(1.9), Inches(0.8), "Theme", (theme_hint or "Auto")[:12], accent_key="accent2")
    add_kpi(slide, Inches(4.8), Inches(5.05), Inches(1.9), Inches(0.8), "Style", "Premium", accent_key="accent3")
    add_textbox(slide, Inches(0.72), Inches(6.35), Inches(5.2), Inches(0.25), f"Created by {PROJECT_NAME}", font_size=10, color_key="accent", bold=True)

    # Content slides
    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bullets = spec.get("bullets") or []
        layout = _content_layout_variant(idx, bullets)
        accent_key = ["accent", "accent2", "accent3"][idx % 3]
        _add_bg(slide, theme, rgb, MSO_SHAPE, Inches, variant=layout)

        slide_title = spec.get("title", f"Slide {idx}")
        kicker = spec.get("kicker") or f"Slide {idx:02d}"

        if layout == "section":
            add_section_header(slide, slide_title, kicker, Inches(0.8), Inches(0.65), Inches(8.8), accent_key=accent_key)
            lead_text = "Use this section to introduce a major idea or transition."
            if bullets:
                first = bullets[0]
                lead_text = first.get("text", "") if isinstance(first, dict) else str(first)
            add_textbox(slide, Inches(0.86), Inches(2.0), Inches(7.6), Inches(1.0), lead_text,
                        font_size=24, color_key="text", bold=True, font_name="Aptos Display")
            if bullets:
                add_bullet_card(slide, Inches(7.55), Inches(1.92), Inches(4.0), Inches(2.6), "Key points", bullets, accent_key=accent_key, dense=True)
            add_textbox(slide, Inches(0.86), Inches(6.05), Inches(7.8), Inches(0.3),
                        spec.get("notes") or "Section divider", font_size=10, color_key="muted")
            continue

        if layout == "split":
            add_section_header(slide, slide_title, kicker, Inches(0.8), Inches(0.62), Inches(8.6), accent_key=accent_key)
            add_bullet_card(slide, Inches(0.82), Inches(1.62), Inches(6.0), Inches(4.95), "Main ideas", bullets or ["Add supporting content."], accent_key=accent_key)
            right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.18), Inches(1.62), Inches(4.9), Inches(4.95))
            _fill_shape(right, theme["surface"], rgb)
            right.fill.transparency = 0.04
            add_textbox(slide, Inches(7.5), Inches(1.92), Inches(4.0), Inches(0.4), "Highlights", font_size=14, color_key=accent_key, bold=True, all_caps=True)
            highlight = bullets[:3] if bullets else ["Use this side for a metric, summary, chart, or visual note."]
            y = 2.35
            for item in highlight:
                txt = item.get("text", "") if isinstance(item, dict) else str(item)
                add_textbox(slide, Inches(7.52), Inches(y), Inches(4.0), Inches(0.48), f"- {txt}", font_size=15, color_key="text")
                y += 0.58
            add_textbox(slide, Inches(7.52), Inches(5.9), Inches(4.0), Inches(0.35), spec.get("notes") or "Visual callout", font_size=10, color_key="muted")
            continue

        add_section_header(slide, slide_title, kicker, Inches(0.8), Inches(0.62), Inches(8.6), accent_key=accent_key)
        add_bullet_card(slide, Inches(0.82), Inches(1.58), Inches(7.2), Inches(4.9), "Overview", bullets or ["Add supporting content."], accent_key=accent_key)
        add_kpi(slide, Inches(8.35), Inches(1.68), Inches(3.05), Inches(0.95), "Status", spec.get("status") or "Ready", accent_key=accent_key)
        add_kpi(slide, Inches(8.35), Inches(2.78), Inches(3.05), Inches(0.95), "Focus", spec.get("focus") or "Clarity", accent_key="accent2")
        add_kpi(slide, Inches(8.35), Inches(3.88), Inches(3.05), Inches(0.95), "Type", spec.get("type") or "Content", accent_key="accent3")
        if spec.get("notes"):
            add_textbox(slide, Inches(8.35), Inches(5.15), Inches(3.0), Inches(0.65), str(spec["notes"]), font_size=10, color_key="muted")

    prs.save(output_path)
    if auto_open:
        _open_file(output_path)
    return f"Presentation created: {output_path}"


def _choose_chart(chart_type: str, WorkbookClasses):
    Workbook, BarChart, LineChart, PieChart, Reference, Alignment, Font, PatternFill, get_column_letter = WorkbookClasses
    t = (chart_type or "bar").lower()
    if t in ("line", "trend"):
        return LineChart()
    if t in ("pie", "doughnut"):
        return PieChart()
    return BarChart()


def create_spreadsheet(parameters: dict, player=None) -> str:
    WorkbookClasses = _import_openpyxl()
    Workbook, BarChart, LineChart, PieChart, Reference, Alignment, Font, PatternFill, get_column_letter = WorkbookClasses

    title = (parameters.get("title") or parameters.get("name") or "Workbook").strip()
    output_path = _resolve_output_path(parameters.get("output_path"), title, ".xlsx")
    auto_open = parameters.get("auto_open", True)
    sheets = _parse_json_arg(parameters.get("worksheets"), None) or _parse_json_arg(parameters.get("sheets"), None)
    if not sheets:
        sheets = [{
            "name": "Sheet1",
            "headers": ["Item", "Value"],
            "rows": [["Example", 1], ["Update the tool input with your real data.", ""]],
        }]

    if isinstance(sheets, dict):
        sheets = [sheets]

    wb = Workbook()
    default_sheet = wb.active
    first_sheet = True

    theme = _office_theme()
    fill_header = PatternFill("solid", fgColor=theme["accent"])
    fill_title = PatternFill("solid", fgColor=theme["bg"])
    white_font = Font(color="FFFFFF", bold=True)
    bold_font = Font(bold=True, color=theme["bg"])
    muted_font = Font(color=theme["muted"])
    center = Alignment(horizontal="center", vertical="center")
    wrap = Alignment(wrap_text=True, vertical="top")

    for sheet_spec in sheets:
        name = _sanitize_filename(sheet_spec.get("name", "Sheet"), "Sheet")[:31]
        ws = default_sheet if first_sheet else wb.create_sheet(title=name)
        ws.title = name
        first_sheet = False

        headers = sheet_spec.get("headers") or []
        rows = sheet_spec.get("rows") or []
        title_row = sheet_spec.get("title")
        start_row = 1

        if title_row:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 2, len(rows[0]) if rows else 2))
            c = ws.cell(1, 1, title_row)
            c.font = Font(bold=True, size=14, color="FFFFFF")
            c.fill = fill_title
            c.alignment = Alignment(horizontal="left", vertical="center")
            start_row = 3

        if headers:
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(start_row, col_idx, header)
                cell.fill = fill_header
                cell.font = white_font
                cell.alignment = center
            data_start = start_row + 1
        else:
            data_start = start_row

        for row_idx, row in enumerate(rows, data_start):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row_idx, col_idx)
                cell.value = _coerce_cell_value(value)
                cell.alignment = wrap if isinstance(value, str) and len(value) > 24 else Alignment(vertical="top")
                if row_idx % 2 == 0 and headers:
                    cell.fill = PatternFill("solid", fgColor="F4FBFF")
                if isinstance(value, (int, float)) and col_idx > 1:
                    cell.number_format = "#,##0.00"

        if headers:
            end_row = data_start + max(len(rows) - 1, 0)
            end_col = len(headers)
            ws.auto_filter.ref = f"{ws.cell(start_row,1).coordinate}:{ws.cell(start_row + max(len(rows), 1), end_col).coordinate}"
            ws.freeze_panes = ws.cell(data_start, 1)
        else:
            end_row = data_start + max(len(rows) - 1, 0)
            end_col = max((len(r) for r in rows), default=1)

        # optional chart
        chart_spec = sheet_spec.get("chart")
        if chart_spec and rows:
            chart = _choose_chart(chart_spec.get("type", "bar"), WorkbookClasses)
            chart.title = chart_spec.get("title", "")
            chart.style = 2
            chart.y_axis.title = chart_spec.get("y_axis", "")
            chart.x_axis.title = chart_spec.get("x_axis", "")

            chart_headers_row = start_row if headers else data_start - 1
            data_start_row = data_start
            data_end_row = data_start + len(rows) - 1
            if headers and len(headers) >= 2:
                data = Reference(ws, min_col=2, min_row=chart_headers_row, max_row=data_end_row)
                cats = Reference(ws, min_col=1, min_row=data_start_row, max_row=data_end_row)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                ws.add_chart(chart, chart_spec.get("anchor", "E2"))

        # sensible widths
        for col_idx in range(1, end_col + 1):
            values = [ws.cell(r, col_idx).value for r in range(1, min(ws.max_row, 200) + 1)]
            max_len = max([len(str(v)) for v in values if v is not None] or [10])
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max(max_len + 2, 12), 36)

        ws.row_dimensions[1].height = 22
        if title_row:
            ws.row_dimensions[1].height = 26

    if "Sheet" in wb.sheetnames and len(wb.sheetnames) > 1:
        wb.remove(wb["Sheet"])

    wb.save(output_path)
    if auto_open:
        _open_file(output_path)
    return f"Spreadsheet created: {output_path}"
